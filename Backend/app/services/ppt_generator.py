import os
from io import BytesIO
from datetime import datetime, date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from sqlalchemy import text
from sqlalchemy.orm import Session
import requests 
from app.models.tank_header import Tank
from app.models.tank_details import TankDetails
from app.models.tank_images_model import TankImages
from app.models.tank_regulations import TankRegulation
from app.models.regulations_master import RegulationsMaster
from app.models.cargo_tank import CargoTankTransaction
from app.models.cargo_master import CargoTankMaster
from app.models.tank_certificate import TankCertificate
from app.models.tank_drawings import TankDrawing
from app.models.valve_test_report import ValveTestReport
from app.models.multiple_regulation import MultipleRegulation
from app.models.product_master_model import ProductMaster
from app.models.safety_valve_brand_model import SafetyValveBrand
from app.models.inspection_valve_model import InspectionValve
from app.models.inspection_gauge_model import InspectionGauge
from app.models.tank_frame_outer_model import TankFrameOuter
from app.models.other_images_model import TankOtherImage
from pathlib import Path
from app.utils.s3_utils import to_cdn_url, s3_client, AWS_S3_BUCKET
# ---------------------------------------------------------------------------
# CONFIGURATION
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).resolve().parent.parent
# This is the folder that contains your "uploads" directory.
# If your uploads live somewhere else, change this accordingly.
MEDIA_ROOT = BASE_DIR  # so MEDIA_ROOT / "uploads/..." matches your DB path

THEME_COLOR = RGBColor(0, 51, 102)
BAHNSCHRIFT_CONDENSED = "Bahnschrift Condensed"
HEADER_TEXT = "Smart-Gas Pte Ltd"
SUB_HEADER_TEXT = (
    "140 Paya Lebar Road #05-21, AZ Building, Singapore 409015\n"
    "Tel : (65) 6848 1040 / Fax : (65) 6848 2173"
)

IMAGE_TYPE_MAP = {
    "frontview": "Front View",
    "rearview": "Rear View",
    "topview": "Top View",
    "undersideview01": "Underside View 01",
    "undersideview02": "Underside View 02",
    "frontlhview": "Front LH View",
    "rearlhview": "Rear LH View",
    "frontrhview": "Front RH View",
    "rearrhview": "Rear RH View",
    "lhsideview": "LH Side View",
    "rhsideview": "RH Side View",
    "valvessectionview": "Valves Section",
    "safetyvalve": "Safety Valve",
    "levelpressuregauge": "Level/Pressure Gauge",
    "vacuumreading": "Vacuum Reading",
}

VALVE_FEATURES = [
    "Top Fill - Gas", "Emergency Valve - A3", "S/V - 1",
    "Iso - Top Fill", "Trycock - 1", "S/V - 2",
    "Bottom fill - Liq", "Trycock - 2 / 3", "S/V - 3",
    "Iso - Bottom Fill - A3", "Vacuum Valve", "S/V - 4",
    "Iso - Top / Bottom", "T-Couple Valve", "S/V - B.Disc",
    "Drain Valve", "T-Couple DV-6", "S/V Diverter",
    "Blow Valve", "Liq Connection", "Line SRV - 1",
    "PB Valve Inlet", "Gas Connection", "Line SRV - 2",
    "PB Valve Return", "Vent Pipe", "Line SRV - 3",
    "PB Unit", "Pipe / Support", "Line SRV - 4 / 5",
    "Sample Valve", "Check Valve", "Out Ves B.Disc"
]

GAUGE_FEATURES = [
    "Pressure Gauge", "Tank Number", "Top",
    "Liquid Gauge", "Co. Logo", "Bottom",
    "Temp Gauge", "TEIP", "Front",
    "Gas Phase Valve", "Haz Ship Label", "Rear",
    "Liquid Phase Valve", "Handling Label", "Left",
    "Equalizing Valve", "Weight Label", "Right",
    "Pump - Smith", "T-75 / 22K7 Label", "Cross Member",
    "Motor", "Tank P&ID Plate", "Cab Door - Rear",
    "Electrical Panel", "Tank Data Plate", "Cab Door Lock",
    "Electrical Plug", "Tank CSC Plate", "Paint Condition",
    "Pump / Motor Mounting", "Std Identification Label", "Reflective Marking"
]



def normalize_image_path(path: str) -> str:
    """Normalize image path for consistent comparison by removing uploads/ prefix and standardizing slashes."""
    if not path:
        return ""
    normalized = str(path).replace("\\", "/").lstrip("/")
    if normalized.startswith("uploads/"):
        normalized = normalized.replace("uploads/", "", 1)
    return normalized


def get_image_key(path: str, tank_number: str = None, base_dir: str = None) -> str:
    """
    Get a unique key for an image that works across all path variants.
    This normalizes DB paths, resolved paths, and CDN keys to the same identifier.
    """
    if not path:
        return ""

    # Convert to forward slashes and strip leading slashes
    key = str(path).replace("\\", "/").lstrip("/")

    # For absolute paths, try to extract the relative part after uploads/
    if key.startswith("uploads/"):
        key = key.replace("uploads/", "", 1)
    elif "uploads/" in key:
        # Handle cases like /path/to/uploads/iso_tank/...
        parts = key.split("uploads/")
        if len(parts) > 1:
            key = parts[1]

    # Remove any remaining absolute path prefixes that might interfere
    # But keep the relative path structure
    return key


def resolve_path(file_path, tank_number, base_dir):
    """
    Try to resolve a DB path like:
      uploads/iso_tank/2025/12/1765332438_smau_8880704_front_view.jpg

    into an actual local filesystem path, checking multiple candidate folders.
    """
    if not file_path:
        return None

    # Standardize slashes
    clean_db_path = str(file_path).replace("\\", "/")

    # Remove 'uploads/' prefix if present so we don't end up with uploads/uploads/...
    if clean_db_path.startswith("uploads/"):
        clean_db_path = clean_db_path.replace("uploads/", "", 1)

    filename_only = os.path.basename(clean_db_path)


    current_uploads = os.path.join(base_dir, "uploads")

    # 2) Another project sibling: ISOTank-Mobile (1)/Backend/uploads
    external_uploads = os.path.abspath(
        os.path.join(
            base_dir,
            ".", ".",                 # stay relative, then go up if needed
            "ISOTank-Mobile (1)",
            "Backend",
            "uploads",
        )
    )

    candidates = [
        # --- Local backend uploads ---
        os.path.join(current_uploads, clean_db_path),
        os.path.join(current_uploads, "tank_images_mobile", clean_db_path),
        os.path.join(current_uploads, "tank_images_mobile", tank_number, "originals", filename_only),

        # --- Other project uploads (mobile backend) ---
        os.path.join(external_uploads, clean_db_path),
        os.path.join(external_uploads, "tank_images_mobile", clean_db_path),
        os.path.join(external_uploads, "tank_images_mobile", tank_number, "originals", filename_only),
        os.path.join(external_uploads, "tank_images_mobile", tank_number, "thumbnails", filename_only),
        os.path.join(external_uploads, "drawings", tank_number, filename_only),
        os.path.join(external_uploads, "certificates", tank_number, filename_only),
    ]

    for path in candidates:
        if os.path.exists(path):
            return path

    return None


def add_picture_smart(slide, path_or_key, left, top, width=None, height=None):
    """
    Adds a picture to the slide, trying local file first, then S3/CDN.
    Returns the picture shape if successful, else None.
    """
    if not path_or_key:
        return None

    # 1) Try Local File
    if os.path.exists(str(path_or_key)):
        try:
            return slide.shapes.add_picture(path_or_key, left, top, width=width, height=height)
        except Exception:
            pass

    fetched_bio = None
    
    # 2) Try S3/CDN
    # Standardize to forward slashes for S3/CDN
    s3_key = str(path_or_key).replace("\\", "/").lstrip("/")

    # Normalize key: ensure it's relative to bucket root (usually starts with 'uploads/')
    # If it's missing 'uploads/' and isn't an absolute local path, try both or prepend.
    if not s3_key.startswith("uploads/") and not os.path.isabs(s3_key):
        s3_key = f"uploads/{s3_key}"

    # A) Try CDN
    try:
        url = to_cdn_url(s3_key)
        if url and (str(url).startswith("http://") or str(url).startswith("https://")):
            r = requests.get(url, timeout=12)
            if r.status_code == 200 and r.content:
                fetched_bio = BytesIO(r.content)
    except Exception:
        pass

    # B) Try Direct S3 Client
    if not fetched_bio:
        try:
            if s3_client and AWS_S3_BUCKET:
                obj = s3_client.get_object(Bucket=AWS_S3_BUCKET, Key=s3_key)
                data = obj.get("Body").read()
                if data:
                    fetched_bio = BytesIO(data)
        except Exception:
            pass

    if fetched_bio:
        try:
            return slide.shapes.add_picture(fetched_bio, left, top, width=width, height=height)
        except Exception:
            pass

    return None


def format_value(value, suffix: str = "") -> str:
    if value is None or value == "":
        return "-"
    if isinstance(value, bool):
        return "Yes" if value else "No"
    if isinstance(value, (datetime, date)):
        return value.strftime("%d-%b-%Y")
    try:
        float_val = float(value)
        if float_val.is_integer():
            return f"{int(float_val):,} {suffix}".strip()
        # NOTE: :,.2f uses comma grouping and 2 decimals
        return f"{float_val:,.2f} {suffix}".strip()
    except Exception:
        return f"{str(value)} {suffix}".strip()


# ---------------------------------------------------------------------------
# LAYOUT HELPERS
# ---------------------------------------------------------------------------

def add_custom_header(slide, title_text: str = "", center_title: bool = False, show_logo: bool = True) -> None:
    """Add the company logo and report title to the header."""
    if show_logo:
        logo_path = os.path.join(BASE_DIR, "logo.jpg")
        if os.path.exists(logo_path):
            try:
                slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.2), height=Inches(0.8))
            except Exception:
                pass

    display_title = "Tank Container Inspection Report" if not title_text else title_text
    
    if center_title:
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tx_box.text_frame.word_wrap = True
        p = tx_box.text_frame.paragraphs[0]
        p.text = display_title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.name = BAHNSCHRIFT_CONDENSED
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    else:
        tx_box = slide.shapes.add_textbox(Inches(4.5), Inches(0.3), Inches(5), Inches(0.6))
        p = tx_box.text_frame.paragraphs[0]
        p.text = display_title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.name = BAHNSCHRIFT_CONDENSED
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT


def add_footer(slide, tank_number: str = "", report_number: str = "", page_num: int = 1):
    """Add a detailed footer as per the reference image."""
    # Horizontal line - Set to 6.5 (Moved up from 6.85 to be inside slide)
    line_y = Inches(6.5) 
    line = slide.shapes.add_connector(
        1,  # straight line
        Inches(0.5), line_y,
        Inches(9.5), line_y
    )
    line.line.color.rgb = THEME_COLOR
    line.line.width = Pt(1)

    # Center: Company Name + Address
    # Set to 6.55 (Moved up from 6.90)
    tx_center = slide.shapes.add_textbox(Inches(3.0), Inches(6.55), Inches(4), Inches(0.8))
    tf = tx_center.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Smart-Gas Pte Ltd"
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "140 Paya Lebar Road #05-21, AZ Building, Singapore 409015"
    p2.font.size = Pt(11)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Tel : (65) 6848 1040 / Fax : (65) 6848 2173"
    p3.font.size = Pt(11)
    p3.alignment = PP_ALIGN.CENTER

    # Left: Report Number - Set to 6.60 (Moved up from 6.95)
    tx_left = slide.shapes.add_textbox(Inches(0.5), Inches(6.60), Inches(3), Inches(0.4))
    p_left = tx_left.text_frame.paragraphs[0]
    p_left.text = report_number if report_number else f"Tank: {tank_number}"
    p_left.font.size = Pt(12) 
    p_left.font.color.rgb = RGBColor(0, 0, 0)
    p_left.font.bold = True

    # Right: Page Number - Set to 6.60 (Moved up from 6.95)
    tx_right = slide.shapes.add_textbox(Inches(8.5), Inches(6.60), Inches(1), Inches(0.4))
    p_right = tx_right.text_frame.paragraphs[0]
    p_right.text = str(page_num)
    p_right.font.size = Pt(14) 
    p_right.font.bold = True
    p_right.alignment = PP_ALIGN.RIGHT


def create_section_title(slide, text_val: str, left, top, width=Inches(9)):
    tx_box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = tx_box.text_frame.paragraphs[0]
    p.text = text_val
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = THEME_COLOR
    return top + Inches(0.3)


def create_compact_table(slide, headers, data_rows, left, top, width, font_size: int = 9):
    """Create a compact table; returns the y-position after the table."""
    if not data_rows:
        return top

    rows = len(data_rows) + 1
    cols = len(headers)
    row_height = Inches(0.25)

    table = slide.shapes.add_table(
        rows, cols, left, top, width, row_height * rows
    ).table

    # Header row
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Data rows
    for r_idx, row_data in enumerate(data_rows):
        for c_idx, value in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)

    return top + (row_height * rows) + Inches(0.3)


def create_kv_block(slide, title: str, data_pairs, left, top, width):
    """Create a 2-column 'label : value' block."""
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = THEME_COLOR

    table_top = top + Inches(0.3)
    rows = len(data_pairs)
    table_height = Inches(0.22) * rows

    table = slide.shapes.add_table(rows, 2, left, table_top, width, table_height).table
    table.columns[0].width = int(width * 0.45)
    table.columns[1].width = int(width * 0.55)

    for i, (label, value) in enumerate(data_pairs):
        cell_lbl = table.cell(i, 0)
        cell_lbl.text = str(label)
        p_lbl = cell_lbl.text_frame.paragraphs[0]
        p_lbl.font.bold = True
        p_lbl.font.size = Pt(9)

        cell_val = table.cell(i, 1)
        cell_val.text = str(value)
        cell_val.text_frame.paragraphs[0].font.size = Pt(9)

    return table_top + table_height + Inches(0.2)


def get_status_text(status_id):
    if status_id == 1:
        return "Y"
    elif status_id == 2:
        return "X"
    elif status_id == 3:
        return "NA"
    return "-"


def add_checklist_footer(slide, top_y, width=Inches(9.0), left_x=Inches(0.5)):
    """
    Adds the specific legend footer for checklist slides.
    Structure: | Refer to Tank P&ID | Satisfactory - Yes - Y | Not Satisfactory - No - X | Not Applicable - NA |
    """
    # 4 columns
    cols = 4
    height = Inches(0.4)
    
    table = slide.shapes.add_table(1, cols, left_x, top_y, width, height).table
    
    # Text and widths
    texts = [
        "Refer to Tank P&ID", 
        "Satisfactory – Yes - Y", 
        "Not Satisfactory – No - X", 
        "Not Applicable - NA"
    ]
    
    # Width distribution (approximate based on text length)
    # Total width 9 inches.
    # 1: 1.8, 2: 2.2, 3: 2.5, 4: 2.5
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    
    for i, text_val in enumerate(texts):
        cell = table.cell(0, i)
        cell.text = text_val
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME_COLOR
        
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)  # Increased from 14



        p.font.name = BAHNSCHRIFT_CONDENSED
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_checklist_slide(prs, tank_number, report_no, page_num, features_list, data_map, col_headers):
    """
    Creates a checklist slide (Slide 2 or 3).
    data_map: dict of { feature_name: status_id }
    col_headers: list of 3 strings for the 3 major columns (e.g. ["Valves", "Valves", "Valves"])
    """
    slide = add_tank_slide(prs, "", tank_number, report_no, page_num, center_header=False, show_logo=True)
    
    # 1. Sub-header Context Bar (Tank No | Report No)
    bar_top = Inches(0.9) # Moved up from 1.0 to shift table higher
    bar_height = Inches(0.35)
    
    # Left Box (Tank No)
    box_left = slide.shapes.add_textbox(Inches(0.5), bar_top, Inches(2.0), bar_height)
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = RGBColor(242, 242, 242) # Light Gray
    p = box_left.text_frame.paragraphs[0]
    p.text = tank_number
    p.font.bold = True
    p.font.size = Pt(18) 
    p.font.name = BAHNSCHRIFT_CONDENSED
    p.alignment = PP_ALIGN.LEFT
    
    # Right Box (Report No)
    box_right = slide.shapes.add_textbox(Inches(7.5), bar_top, Inches(2.0), bar_height)
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = RGBColor(242, 242, 242)
    p = box_right.text_frame.paragraphs[0]
    p.text = report_no
    p.font.bold = True
    p.font.size = Pt(18) 
    p.font.name = BAHNSCHRIFT_CONDENSED
    p.alignment = PP_ALIGN.RIGHT

    # 2. Checklist Table
    table_top = Inches(1.35) # Moved up from 1.45
    table_width = Inches(9.0)
    
    import math
    num_cols = 3
    num_items = len(features_list)
    num_rows = math.ceil(num_items / num_cols)
    
    total_rows = num_rows + 1 # +1 for header
    total_cols = 6 # 3 pairs of (Label, Status)
    
    row_height = Inches(0.4)
    
    tbl = slide.shapes.add_table(total_rows, total_cols, Inches(0.5), table_top, table_width, row_height * total_rows).table
    
    cw_label = Inches(2.3)
    cw_status = Inches(0.7)
    
    for i in range(3):
        tbl.columns[i*2].width = cw_label
        tbl.columns[i*2+1].width = cw_status
        
    # Header Row
    for i, header_text in enumerate(col_headers):
        c1 = tbl.cell(0, i*2)
        c2 = tbl.cell(0, i*2+1)
        c1.merge(c2)
        
        c1.text = header_text
        c1.fill.solid()
        c1.fill.fore_color.rgb = THEME_COLOR
        p = c1.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.name = BAHNSCHRIFT_CONDENSED
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(19) # Increased from 16



    # Data Rows
    for r in range(num_rows):
        for c in range(3):
            idx = r * 3 + c
            if idx < len(features_list):
                feature = features_list[idx]
                status_id = data_map.get(feature)
                status_text = get_status_text(status_id)
                
                # Label Cell
                c_lbl = tbl.cell(r+1, c*2)
                c_lbl.text = feature
                c_lbl.fill.solid()
                c_lbl.fill.fore_color.rgb = RGBColor(230, 230, 230) 
                p = c_lbl.text_frame.paragraphs[0]
                p.font.size = Pt(16) # Increased from 13
                p.font.name = BAHNSCHRIFT_CONDENSED
                p.font.color.rgb = RGBColor(0, 0, 0)
                
                # Status Cell
                c_st = tbl.cell(r+1, c*2+1)
                c_st.text = status_text
                c_st.fill.solid()
                c_st.fill.fore_color.rgb = RGBColor(242, 242, 242)
                p = c_st.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(18) # Increased from 15
                p.font.bold = True
                p.font.name = BAHNSCHRIFT_CONDENSED
                
    # 3. Legend Footer
    legend_top = table_top + (row_height * total_rows) + Inches(0.1)
    add_checklist_footer(slide, legend_top)




def add_tank_slide(prs: Presentation, title_text: str = "", tank_number: str = "", report_number: str = "", page_num: int = 1, center_header: bool = False, show_logo: bool = True):
    """Create a blank slide with custom header and footer."""
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name and layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank_layout)
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            el = shape._element
            el.getparent().remove(el)

    add_custom_header(slide, title_text, center_header, show_logo)
    add_footer(slide, tank_number, report_number, page_num)
    return slide


def add_image_bar(slide, text_val: str, left, top, width):
    """Add a dark blue bar with white text."""
    bar = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, Inches(0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME_COLOR
    bar.line.fill.background() # No border
    
    tf = bar.text_frame
    tf.text = text_val
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_image_with_bars(slide, path, left, top, width, height, label):
    """Add image with blue bars above and below it."""
    # Top bar
    add_image_bar(slide, label, left, top, width)
    
    # Image
    pic = add_picture_smart(slide, path, left, top + Inches(0.35), width=width, height=height)
    if not pic and path:
        # If we failed to render an existing path, show error text
        err_box = slide.shapes.add_textbox(left, top + Inches(0.35), width, height)
        err_box.text_frame.text = "[Image Not Found / Error Rendering]"
        err_p = err_box.text_frame.paragraphs[0]
        err_p.alignment = PP_ALIGN.CENTER
        err_p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Bottom bar
    add_image_bar(slide, label, left, top + height + Inches(0.4), width)


# ---------------------------------------------------------------------------
# IMAGE SEQUENCES
# ---------------------------------------------------------------------------

def add_image_sequence(prs: Presentation, tank_number: str, image_list, section_title: str, base_dir: str, exclude_paths: set = None) -> None:
    """
    Adds one slide per image in image_list.
    Each item in image_list is expected to be:
       { "path": <db_path>, "label": <text> }
    exclude_paths: A set of file paths (raw or resolved) to skip if encountered.
    """
    valid_items = []
    if exclude_paths is None:
        exclude_paths = set()

    for item in image_list:
        raw_path = item.get("path")
        if not raw_path:
            continue

        # Check raw path exclusion
        if raw_path in exclude_paths:
            continue

        real_path = resolve_path(raw_path, tank_number, base_dir)
        
        # Check resolved path exclusion
        if real_path and real_path in exclude_paths:
            continue

        if real_path:
            valid_items.append(
                {"path": real_path, "label": item.get("label", "")}
            )
            exclude_paths.add(real_path)
        else:
            # Keep original path for debugging text
            valid_items.append(
                {
                    "path": None,
                    "label": item.get("label", ""),
                    "original": raw_path,
                }
            )
        
        # Always exclude raw path too once used
        exclude_paths.add(raw_path)

    if not valid_items:
        return

    for img_item in valid_items:
        path = img_item.get("path")
        label = img_item.get("label", section_title)
        
        # Use center header style (True) and hide logo (False) to match tank condition slides
        page_num = prs.slides.index(prs.slides[-1]) + 2
        current_slide = add_tank_slide(prs, f"{tank_number} – {section_title}", tank_number, "", page_num, True, False)

        img_width = Inches(9.0)
        img_height = Inches(4.5)
        img_left = Inches(0.5)
        img_top = Inches(1.15)

        # We can reuse add_image_with_bars but we need to handle the S3/CDN fallback inside add_image_sequence
        # if add_picture_smart fails or we don't have a local path.
        
        # Top bar
        add_image_bar(current_slide, label, img_left, img_top, img_width)
        
        # Image logic (consolidated with add_picture_smart logic)
        path_to_use = path if path else img_item.get("original")
        pic = add_picture_smart(current_slide, path_to_use, img_left, img_top + Inches(0.35), width=img_width, height=img_height)
        
        if not pic:
            err_box = current_slide.shapes.add_textbox(img_left, img_top + Inches(0.35), img_width, img_height)
            err_box.text_frame.text = f"[Image Not Found / Error Rendering]\n{path_to_use or ''}"
            err_p = err_box.text_frame.paragraphs[0]
            err_p.alignment = PP_ALIGN.CENTER
            err_p.font.color.rgb = RGBColor(255, 0, 0)
        else:
            # Re-center if it was smaller than full width
            pic.left = int((prs.slide_width - pic.width) / 2)

        # Bottom bar
        add_image_bar(current_slide, label, img_left, img_top + img_height + Inches(0.4), img_width)


def add_full_page_image_slide(prs, tank_number, report_no, page_num, slide_title, top_bar_text, bottom_bar_text, image_path, base_dir, exclude_paths=None):
    """Adds a slide with a single large image, a slide title, and top/bottom blue bars."""
    if not image_path:
        return page_num

    # Get unified key for this image
    image_key = get_image_key(image_path, tank_number, base_dir)

    # Check if this image has already been used
    if exclude_paths is not None and image_key in exclude_paths:
        return page_num

    real_path = resolve_path(image_path, tank_number, base_dir)
    # If path resolving failed, we might still try to render if it's a valid remote path, 
    # but strictly checking real_path usually avoids missing file errors. 
    # For now, if no real_path found, we skip to avoid empty slides (unless it's an S3 path).
    # But resolve_path checks S3/local candidates. 
    # Let's rely on add_picture_smart to handle final resolution if resolve_path fails but path exists.

    # Check resolved path key as well
    if real_path:
        resolved_key = get_image_key(real_path, tank_number, base_dir)
        if exclude_paths is not None and resolved_key in exclude_paths:
            return page_num
    
    # Add to seen paths
    if exclude_paths is not None:
        exclude_paths.add(image_key)
        if real_path:
            exclude_paths.add(get_image_key(real_path, tank_number, base_dir))

    # Create slide
    # Slide Title is passed explicitly.
    slide = add_tank_slide(prs, slide_title, tank_number, report_no, page_num, center_header=True, show_logo=False)

    img_left = Inches(0.5)
    img_top = Inches(1.15)
    img_width = Inches(9.0)
    img_height = Inches(4.5)

    # Top Bar
    if top_bar_text:
        add_image_bar(slide, top_bar_text, img_left, img_top, img_width)

    # Image
    pic_top = img_top + Inches(0.35)
    pic = add_picture_smart(slide, image_path, img_left, pic_top, width=img_width, height=img_height)
    
    # Error handling / Centering
    if not pic:
         # If primary image failed, maybe try a placeholder or just leave empty?
         # Current logic: show error text if desired, or nothing.
         pass
    else:
        # Re-center if it was smaller than full width
        pic.left = int((prs.slide_width - pic.width) / 2)

    # Bottom Bar
    if bottom_bar_text:
        bar_top = img_top + img_height + Inches(0.4)
        add_image_bar(slide, bottom_bar_text, img_left, bar_top, img_width)

    return page_num + 1


def add_images_grid(prs, tank_number, report_no, start_page_num, title, image_items, base_dir, exclude_paths=None):
    """
    Adds slides displaying images in a 3x2 grid.
    image_items: list of {"path": str, "label": str}
    """
    if not image_items:
        return start_page_num

    curr_page = start_page_num
    
    # Filter
    valid_items = []
    if exclude_paths is None:
        exclude_paths = set()
        
    for item in image_items:
        raw = item.get("path")
        if not raw:
            continue

        # Get unified key for this image
        image_key = get_image_key(raw, tank_number, base_dir)

        # Check if this image has already been used
        if image_key in exclude_paths:
            continue

        resolved = resolve_path(raw, tank_number, base_dir)
        if resolved:
            resolved_key = get_image_key(resolved, tank_number, base_dir)
            if resolved_key in exclude_paths:
                continue

        # Image is available, add it
        valid_items.append({"path": raw, "resolved_path": resolved, "label": item.get("label")})
        # Mark as used
        exclude_paths.add(image_key)
        if resolved:
            exclude_paths.add(resolved_key)
            
    if not valid_items:
        return curr_page

    # 3x2 Grid
    chunk_size = 6
    for i in range(0, len(valid_items), chunk_size):
        chunk = valid_items[i:i+chunk_size]
        
        slide = add_tank_slide(prs, title, tank_number, report_no, curr_page, center_header=True, show_logo=False)
        
        # Grid settings
        # Start below header (approx 1.2 inch)
        start_y = Inches(1.3)
        # Margins
        margin_x = Inches(0.5)
        
        # 3 columns in 9 inches -> ~2.8 inch width per image with gaps
        # Gap 0.3 -> 2.8*3 + 0.3*2 = 8.4+0.6 = 9.0. Good.
        col_w = Inches(2.8)
        gap_x = Inches(0.3)
        
        # 2 rows in ~5.5 inches -> ~2.5 inch height
        row_h = Inches(2.5)
        gap_y = Inches(0.3)
        
        for j, img in enumerate(chunk):
            row = j // 3
            col = j % 3
            
            x = margin_x + (col * (col_w + gap_x))
            y = start_y + (row * (row_h + gap_y))
            
            # Add picture
            # We want to fit it in the box - Force fixed dimensions to prevent layout shifts
            try:
                # Force both width and height to fix the grid layout
                add_picture_smart(slide, img["path"], x, y, width=col_w, height=row_h)
            except Exception:
                pass
                
        curr_page += 1
        
    return curr_page


def add_contact_sheet_slide(prs, tank_number, report_no, start_page_num, image_items, base_dir, exclude_paths=None):
    """
    Adds a slide (or slides) displaying images in an adaptive grid layout.
    Maximizes image density based on count (1-9). No header bar.
    image_items: list of {"path": str, "label": str}
    """
    if not image_items:
        return start_page_num

    curr_page = start_page_num
    
    # Filter valid items
    valid_items = []
    if exclude_paths is None:
        exclude_paths = set()
        
    for item in image_items:
        raw = item.get("path")
        if not raw:
            continue
        valid_items.append(item)
    
    if not valid_items:
        return curr_page

    chunk_size = 9
    margin_x = Inches(0.5)
    margin_y = Inches(1.2) # Increased from 0.5 to 1.2 to clear header
    
    # Process in chunks of 9
    for i in range(0, len(valid_items), chunk_size):
        chunk = valid_items[i:i+chunk_size]
        n_imgs = len(chunk)
        
        # Use simple blank slide with footer
        slide = add_tank_slide(prs, "", tank_number, report_no, curr_page, center_header=False, show_logo=False)
        
        # Determine Grid Dimensions dynamically based on n_imgs
        if n_imgs == 1:
            rows, cols = 1, 1
        elif n_imgs == 2:
            rows, cols = 1, 2 # Side by side
        elif n_imgs <= 4:
            rows, cols = 2, 2
        elif n_imgs <= 6:
            rows, cols = 2, 3
        else:
            rows, cols = 3, 3

        # Calculate cell sizes
        # Usable area: W=9.0, H=5.3 (approx) - Reduced height to accommodate larger top margin
        usable_w = Inches(9.0)
        usable_h = Inches(5.3) 
        
        cell_w = usable_w / cols
        cell_h = usable_h / rows
        
        for j, img_item in enumerate(chunk):
            row = j // cols
            col = j % cols
            
            # Base position
            base_x = margin_x + (col * cell_w)
            base_y = margin_y + (row * cell_h)
            
            path = img_item.get("path")
            # Resolve path if needed
            real_path = resolve_path(path, tank_number, base_dir)
            
            try:
                # Add picture forcing FIXED width and height to fill the cell
                # This ensures the layout does not change "as per images" (aspect ratios)
                add_picture_smart(slide, real_path or path, base_x, base_y, width=cell_w, height=cell_h)
            except Exception:
                pass

        curr_page += 1
        
    return curr_page


# ---------------------------------------------------------------------------
# MAIN GENERATOR
# ---------------------------------------------------------------------------

def resolve_product_name(db: Session, insp: dict | None, d: TankDetails | None, cargos) -> str:
    """
    Resolve product name with priority:
    1) inspection.product / product_id
    2) TankDetails.product_id
    3) last cargo product
    """
    # --- 1. Inspection-based ---
    if insp:
        # Priority: 1) Joined name 2) ID lookup 3) Raw value
        joined_name = insp.get("product_name")
        if joined_name and joined_name not in (None, "", "-"):
            return str(joined_name)

        raw = insp.get("product_id") or insp.get("product")
        if raw not in (None, "", "-"):
            try:
                if str(raw).isdigit():
                    prod = (
                        db.query(ProductMaster)
                        .filter(ProductMaster.id == int(raw))
                        .first()
                    )
                    if prod:
                        return prod.product_name
                else:
                    return str(raw)
            except Exception:
                pass

    # --- 2. TankDetails-based ---
    if d and d.product_id:
        try:
            prod = (
                db.query(ProductMaster)
                .filter(ProductMaster.id == d.product_id)
                .first()
            )
            if prod:
                return prod.product_name
        except Exception:
            pass

    # --- 3. Cargo-based (LAST CARGO – reference behaviour) ---
    if cargos:
        _, cargo_master = cargos[-1]
        if hasattr(cargo_master, "product_name") and cargo_master.product_name:
            return cargo_master.product_name

    return "-"


def resolve_safety_valve_name(db: Session, insp: dict | None, d: TankDetails | None) -> str:
    """
    Resolve safety valve brand name with priority:
    1) inspection joined name
    2) TankDetails.safety_valve_brand_id lookup
    """
    # 1. Check inspection joined name (if it was joined)
    if insp:
        joined = insp.get("safety_valve_brand_name")
        if joined and joined not in (None, "", "-"):
            return str(joined)

    # 2. Check TankDetails fallback
    if d and d.safety_valve_brand_id:
        try:
            brand = (
                db.query(SafetyValveBrand)
                .filter(SafetyValveBrand.id == d.safety_valve_brand_id)
                .first()
            )
            if brand:
                return brand.brand_name
        except Exception:
            pass

    # 3. Check inspection raw id
    if insp:
        raw = insp.get("safety_valve_brand_id") or insp.get("safety_valve_brand")
        if raw and raw not in (None, "", "-"):
            return str(raw)

    return "-"


def create_presentation(db: Session, tank_id: int, base_dir: str | None = None, inspection_id: int | None = None) -> BytesIO:
    """
    Build the PowerPoint presentation for a given tank_id and return a BytesIO.
    If inspection_id is provided, generate for that specific inspection.
    Otherwise, generate for the latest inspection.
    """
    # Default base_dir → project root (two levels above this file)
    if not base_dir:
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # --- Core tank + details ---
    tank = db.query(Tank).filter(Tank.id == tank_id).first()
    d = db.query(TankDetails).filter(TankDetails.tank_id == tank_id).first()

    if not tank or not d:
        raise ValueError("Tank Details not found")

    # Regulations (from multiple_regulation)
    regs = (
        db.query(MultipleRegulation)
        .filter(MultipleRegulation.tank_id == tank_id)
        .all()
    )

    # Other info
    cargos = (
        db.query(CargoTankTransaction, CargoTankMaster)
        .join(CargoTankMaster, CargoTankTransaction.cargo_reference == CargoTankMaster.id)
        .filter(CargoTankTransaction.tank_id == tank_id)
        .all()
    )
    certs = db.query(TankCertificate).filter(TankCertificate.tank_id == tank_id).order_by(TankCertificate.created_at.desc()).all()
    drawings = []
    if d and d.pid_id:
        dwg = db.query(TankDrawing).filter(TankDrawing.id == d.pid_id).first()
        if dwg:
            drawings.append(dwg)
            
    valves = db.query(ValveTestReport).filter(ValveTestReport.tank_id == tank_id).all()
    
    # New Image Data Sources
    tank_frame_outer_record = None
    if d and d.ga_id:
        tank_frame_outer_record = db.query(TankFrameOuter).filter(TankFrameOuter.id == d.ga_id).first()
    
    others_rows = db.query(TankOtherImage).filter(TankOtherImage.tank_id == d.id).all()
    
    # Checklist Data
    insp_valves = db.query(InspectionValve).filter(InspectionValve.tank_id == tank_id).all()
    insp_gauges = db.query(InspectionGauge).filter(InspectionGauge.tank_id == tank_id).all()
    
    valve_map = {v.features: v.status_id for v in insp_valves}
    gauge_map = {g.features: g.status_id for g in insp_gauges}

    # Latest inspection (by tank_id or tank_number) OR specific inspection_id
    if inspection_id:
        insp_sql = text(
            """
            SELECT 
                ti.*, 
                it.inspection_type_name
            FROM tank_inspection_details ti
            LEFT JOIN inspection_type it ON ti.inspection_type_id = it.id
            WHERE ti.inspection_id = :iid AND ti.tank_id = :tid
            """
        )
        params = {"iid": inspection_id, "tid": tank_id}
    else:
        insp_sql = text(
            """
            SELECT 
                ti.*, 
                it.inspection_type_name
            FROM tank_inspection_details ti
            LEFT JOIN inspection_type it ON ti.inspection_type_id = it.id
            WHERE ti.tank_id = :tid OR ti.tank_number = :tn
            ORDER BY ti.inspection_date DESC
            LIMIT 1
            """
        )
        params = {"tid": tank_id, "tn": tank.tank_number}

    insp_row = db.execute(insp_sql, params).mappings().first()
    insp = dict(insp_row) if insp_row else None

    checklist_rows = []
    todo_rows = []
    insp_images = []

    insp_image_map = {}

    if insp:
        iid = insp["inspection_id"]

        # Checklist
        try:
            sql_chk = text(
                """
                SELECT ic.job_name, ic.sub_job_description, COALESCE(ic.status, ist.status_name, '-') as status, ic.comment
                FROM inspection_checklist ic
                LEFT JOIN inspection_status ist ON ic.status_id = ist.status_id
                WHERE ic.inspection_id = :iid
                ORDER BY ic.sub_job_id ASC
                """
            )
            checklist_data = db.execute(sql_chk, {"iid": iid}).fetchall()
            for r in checklist_data:
                checklist_rows.append(
                    [r[0], r[1], r[2] or "-", r[3] or "-"]
                )
        except Exception:
            pass

        # To-do list (faulty items)
        try:
            sql_todo = text(
                """
                SELECT t.job_name, t.sub_job_description, COALESCE(ist.status_name, 'Faulty') as status, t.comment
                FROM to_do_list t
                LEFT JOIN inspection_status ist ON t.status_id = ist.status_id
                WHERE t.inspection_id = :iid
                ORDER BY t.created_at ASC
                """
            )
            todo_data = db.execute(sql_todo, {"iid": iid}).fetchall()
            for r in todo_data:
                todo_rows.append(
                    [r[0], r[1], r[2] or "Faulty", r[3] or "-"]
                )
        except Exception:
            pass

        # Inspection images (tank_images)
        try:
            img_sql = text(
                """
                SELECT image_type, image_path
                FROM tank_images
                WHERE inspection_id = :iid
                ORDER BY id ASC
                """
            )
            img_data = db.execute(img_sql, {"iid": iid}).fetchall()
            for r in img_data:
                # Normalize type slug: lowercase, no spaces, no underscores
                itype = str(r[0]).lower().replace(" ", "").replace("_", "")
                real_path = resolve_path(r[1], tank.tank_number, base_dir)
                # Store local path if found, else original DB path (potential S3 key)
                path_to_store = real_path if real_path else r[1]
                insp_image_map.setdefault(itype, []).append(path_to_store)
                
                readable_label = IMAGE_TYPE_MAP.get(itype, itype.title())
                insp_images.append({"path": r[1], "label": readable_label})
        except Exception:
            pass

    # Fallback logic REMOVED to prevent shadowing of specific images
    print("DEBUG: create_presentation started. insp_image_map keys:", list(insp_image_map.keys()))

    # -----------------------------------------------------------------------
    # PPT GENERATION
    # -----------------------------------------------------------------------

    prs = Presentation()

    # --- SLIDE 1: MAIN REPORT FORM ---
    report_no = insp["report_number"] if insp else "-"
    slide1 = add_tank_slide(prs, "", tank.tank_number, report_no, 1)

    # Colors (Match Reference Image)
    COLOR_HEADER_BG = RGBColor(217, 217, 217) # D9D9D9
    COLOR_LABEL_BG = RGBColor(217, 217, 217)  # D9D9D9
    COLOR_VALUE_BG = RGBColor(242, 242, 242)  # F2F2F2
    FONT_NAME = BAHNSCHRIFT_CONDENSED

    # -------------------------------
    # Inspection Type Header (SIMPLIFIED)
    # -------------------------------
    # -------------------------------
    # Inspection Type Header (SIMPLIFIED - REVERTED)
    # -------------------------------
    # Reverting to the gray bar with simple text as requested: 
    # "please revert back inspection type on how it looked before"
    
    hdr_table = slide1.shapes.add_table(
        1, 1, Inches(0.5), Inches(1.10), Inches(9.0), Inches(0.35) # Reduced height from 0.45 to 0.35
    ).table

    insp_type_text = (insp.get("inspection_type_name") or "Condition Check").upper()
    
    cell_hdr = hdr_table.cell(0, 0)
    cell_hdr.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell_hdr.text = f"INSPECTION TYPE: {insp_type_text}"
    cell_hdr.fill.solid()
    cell_hdr.fill.fore_color.rgb = COLOR_HEADER_BG
    p = cell_hdr.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14) # Reduced font size slightly to match smaller header

    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # -------------------------------
    # Main Data Table
    # -------------------------------
    table_left = Inches(0.5)
    table_top = Inches(1.50) # Moved up slightly from 1.55 to fit better
    table_width = Inches(9.0)
    
    # 14 rows: 
    # 1. Date...
    # ...
    # 13. Applicable Regulation (Merged)
    # 14. Cargo Reference (Merged)
    rows = 14
    cols = 4
    # Total height approx 5.5 inches left. 14 rows -> 0.38 per row is too tall (overlaps footer).
    # Reducing to 0.35. Total height = 4.9 inches. Bottom = 1.5 + 4.9 = 6.4. Plenty of room for footer at 6.85.
    row_height = Inches(0.35)
    
    tbl_shape = slide1.shapes.add_table(rows, cols, table_left, table_top, table_width, row_height * rows)
    table = tbl_shape.table
    
    # Column Widths
    cw_label = Inches(2.0)
    cw_value = Inches(2.5)
    table.columns[0].width = cw_label
    table.columns[1].width = cw_value
    table.columns[2].width = cw_label
    table.columns[3].width = cw_value

    def set_cell(r, c, label, value):
        cell_lbl = table.cell(r, c)
        cell_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_lbl.text = str(label)
        cell_lbl.fill.solid()
        cell_lbl.fill.fore_color.rgb = COLOR_LABEL_BG
        p_lbl = cell_lbl.text_frame.paragraphs[0]
        p_lbl.font.bold = True
        p_lbl.font.size = Pt(12) # Reduced slightly to fit
        p_lbl.font.name = FONT_NAME
        p_lbl.font.color.rgb = RGBColor(0, 0, 0)

        cell_val = table.cell(r, c + 1)
        cell_val.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_val.text = str(value)
        cell_val.fill.solid()
        cell_val.fill.fore_color.rgb = COLOR_VALUE_BG
        p_val = cell_val.text_frame.paragraphs[0]
        p_val.font.bold = True
        p_val.font.size = Pt(14) 
        p_val.font.name = FONT_NAME
        p_val.font.color.rgb = RGBColor(0, 0, 0)

    # -------------------------------
    # Resolve Product & Safety Valve
    # -------------------------------
    product_name = resolve_product_name(db, insp, d, cargos)
    safety_valve_name = resolve_safety_valve_name(db, insp, d)

    # Convert dates
    date_formatted = format_value(insp["inspection_date"]) if insp else "-"
    next_insp = format_value(certs[0].next_insp_date) if certs else "-"
    insp_25 = format_value(certs[0].insp_2_5y_date) if certs else "-"
    cert_no = certs[0].certificate_number if certs else "-"

    # Row 0
    date_val = format_value(insp["inspection_date"]) if insp else "-"
    set_cell(0, 0, "Date", date_val)
    set_cell(0, 2, "Report Number", report_no)

    # Row 1
    # Check if ownership name is in details or we need to join?
    # d.ownership is string name in revised model
    set_cell(1, 0, "Tank No", tank.tank_number)
    set_cell(1, 2, "Ownership", d.ownership or "-")

    # Row 2
    set_cell(2, 0, "YOM", format_value(d.date_mfg))
    set_cell(2, 2, "2.5Y / 5.0Y Insp", insp_25)

    # Row 3
    set_cell(3, 0, "Mfgr / Design", d.mfgr or "-")
    set_cell(3, 2, "Next Insp Date", next_insp)

    # Row 4
    set_cell(4, 0, "Mawp", format_value(d.mawp))
    set_cell(4, 2, "Tank Certificate", cert_no)

    # Row 5 (NEW)
    # Ves Mat logic: Reference shows "SS : -196C to 65C".
    # We might map this from design temp?
    ves_mat_val = "-"
    if d.design_temperature:
         ves_mat_val = f"SS : {d.design_temperature}" # Approximation
    set_cell(5, 0, "Ves Mat", ves_mat_val)
    set_cell(5, 2, "Color – Body / Frame", d.color_body_frame or "-")

    # Row 6
    # Tank UN / ISO Code -> UN Code from tank_details, ISO Code from tank_details
    iso_val = d.tank_iso_code if d.tank_iso_code else "-"
    un_val = d.un_code if d.un_code else "-"
    set_cell(6, 0, "Tank/ ISO Code", f"{iso_val}")
    
    # Vacuum Reading -> from tank_inspection_details (insp dict)
    vac_val = insp.get("vacuum_reading") if insp else "-"
    vac_uom = insp.get("vacuum_uom") if insp else ""
    vac_display = f"{vac_val} {vac_uom}".strip() if vac_val and vac_val != "-" else vac_val or "-"
    set_cell(6, 2, "Vacuum Reading", vac_display)

    # Row 7
    set_cell(7, 0, "Tare Kg", format_value(d.tare_weight_kg))
    set_cell(7, 2, "Safety Valve", safety_valve_name)

    # Row 8
    set_cell(8, 0, "Gross Kg", format_value((d.mpl_kg or 0) + (d.tare_weight_kg or 0)))
    set_cell(8, 2, "Evacuation Valve", d.evacuation_valve or "-")

    # Row 9
    set_cell(9, 0, "Net Kg", format_value(d.mpl_kg))
    set_cell(9, 2, "Product - Last", product_name)

    # Row 10
    set_cell(10, 0, "Capacity Liter", format_value(d.capacity_l))
    check_date = tank.created_at.strftime("%d-%m-%y") if tank.created_at else "-"
    set_cell(10, 2, "Check By / Date", f"{d.created_by or '-'} / {check_date}")

    # Row 11
    set_cell(11, 0, "Pump Type", format_value(d.pump_type))
    # Approved by - User wants this field.
    # We can try updated_by
    appr_date = tank.updated_at.strftime("%d-%m-%y") if tank.updated_at else "-"
    set_cell(11, 2, "Approved By / Date", f"{d.updated_by or '-'} / {appr_date}")

    # Row 12 - Applicable Regulation (Merged)
    reg_names = [r.regulation_name for r in regs if r.regulation_name]
    reg_str = ", ".join(reg_names) if reg_names else "-"
    
    # We already created blank cells in add_table, now merge and fill
    c_reg_lbl = table.cell(12, 0)
    c_reg_lbl.text = "Applicable Regulation"
    c_reg_lbl.fill.solid()
    c_reg_lbl.fill.fore_color.rgb = COLOR_LABEL_BG
    p = c_reg_lbl.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0,0,0)

    c_reg_val = table.cell(12, 1)
    c_reg_val.merge(table.cell(12, 3))
    c_reg_val.text = reg_str
    c_reg_val.fill.solid()
    c_reg_val.fill.fore_color.rgb = COLOR_VALUE_BG
    p = c_reg_val.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0,0,0)

    # Row 13 - Cargo Reference (Merged) showing UN Codes
    c_cargo_lbl = table.cell(13, 0)
    c_cargo_lbl.text = "UN Code"
    c_cargo_lbl.fill.solid()
    c_cargo_lbl.fill.fore_color.rgb = COLOR_LABEL_BG
    p = c_cargo_lbl.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0,0,0)

    # Use d.un_code for Cargo Reference (UN codes)
    cargo_ref_val = d.un_code if d.un_code else "-"
    
    c_cargo_val = table.cell(13, 1)
    c_cargo_val.merge(table.cell(13, 3))
    c_cargo_val.text = cargo_ref_val
    c_cargo_val.fill.solid()
    c_cargo_val.fill.fore_color.rgb = COLOR_VALUE_BG
    p = c_cargo_val.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.name = FONT_NAME
    p.font.color.rgb = RGBColor(0,0,0)





    # --- SLIDES 2 & 3: CHECKLISTS ---
    curr_page = 2
    
    # Slide 2: Valves Checklist
    # Use features directly from the database records
    db_valve_features = [v.features for v in insp_valves]
    valve_map = {v.features: v.status_id for v in insp_valves}

    add_checklist_slide(
        prs, tank.tank_number, report_no, curr_page,
        db_valve_features, valve_map, ["Valves", "Valves", "Valves"]
    )
    curr_page += 1
    
    # Slide 3: Gauges Checklist
    # Use features directly from the database records
    db_gauge_features = [g.features for g in insp_gauges]
    gauge_map = {g.features: g.status_id for g in insp_gauges}

    add_checklist_slide(
        prs, tank.tank_number, report_no, curr_page,
        db_gauge_features, gauge_map, ["Gauge / Pump", "Label / Decal", "Body / Frame"]
    )
    curr_page += 1

    # --- TANK CONDITION (Slides 4+) ---
    cond_header = f"{tank.tank_number} – Tank Condition"

    # Front View + Rear View
    s_fr = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    f_path = insp_image_map.get("frontview", [None])[0]
    r_path = insp_image_map.get("rearview", [None])[0]
    add_image_with_bars(s_fr, f_path, Inches(0.5), Inches(1.15), Inches(4.2), Inches(3.5), "Front View")
    add_image_with_bars(s_fr, r_path, Inches(5.3), Inches(1.15), Inches(4.2), Inches(3.5), "Rear View")

    # Top View + Underside Views
    s_tu = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    t_path = insp_image_map.get("topview", [None])[0]
    u1_path = insp_image_map.get("undersideview01", [None])[0]
    u2_path = insp_image_map.get("undersideview02", [None])[0]
    
    # Left Column: Top View Large
    add_image_bar(s_tu, "Top View", Inches(0.5), Inches(1.15), Inches(4.2))
    add_picture_smart(s_tu, t_path, Inches(0.5), Inches(1.50), width=Inches(4.2), height=Inches(4.55))
    add_image_bar(s_tu, "Top View", Inches(0.5), Inches(6.1), Inches(4.2))
    
    # Right Column: Underside Stacked
    add_image_bar(s_tu, "Underside View", Inches(5.3), Inches(1.15), Inches(4.2))
    add_picture_smart(s_tu, u1_path, Inches(5.3), Inches(1.50), width=Inches(4.2), height=Inches(2.1))
    add_picture_smart(s_tu, u2_path, Inches(5.3), Inches(3.95), width=Inches(4.2), height=Inches(2.1))
    add_image_bar(s_tu, "Underside View", Inches(5.3), Inches(6.1), Inches(4.2))

    # Front LH View + Rear LH View
    s_lh = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    flh = insp_image_map.get("frontlhview", [None])[0]
    rlh = insp_image_map.get("rearlhview", [None])[0]
    add_image_with_bars(s_lh, flh, Inches(0.5), Inches(1.15), Inches(4.2), Inches(3.5), "Front LH View")
    add_image_with_bars(s_lh, rlh, Inches(5.3), Inches(1.15), Inches(4.2), Inches(3.5), "Rear LH View")

    # Front RH View + Rear RH View
    s_rh = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    frh = insp_image_map.get("frontrhview", [None])[0]
    rrh = insp_image_map.get("rearrhview", [None])[0]
    add_image_with_bars(s_rh, frh, Inches(0.5), Inches(1.15), Inches(4.2), Inches(3.5), "Front RH View")
    add_image_with_bars(s_rh, rrh, Inches(5.3), Inches(1.15), Inches(4.2), Inches(3.5), "Rear RH View")

    # LH Side View (Single)
    s_lhs = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    lhsv = insp_image_map.get("lhsideview", [None])[0]
    add_image_with_bars(s_lhs, lhsv, Inches(0.5), Inches(1.15), Inches(9.0), Inches(4.5), "LH Side View")

    # RH Side View (Single)
    s_rhs = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    rhsv = insp_image_map.get("rhsideview", [None])[0]
    add_image_with_bars(s_rhs, rhsv, Inches(0.5), Inches(1.15), Inches(9.0), Inches(4.5), "RH Side View")

    # Valves Section View
    s_vals = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    valves_list_imgs = insp_image_map.get("valvessectionview", [])
    v_top = Inches(1.15)
    add_image_bar(s_vals, "Valves Section View", Inches(0.5), v_top, Inches(9.0))
    if valves_list_imgs:
        v_h = Inches(4.5) / len(valves_list_imgs)
        for i, v_path in enumerate(valves_list_imgs):
            add_picture_smart(s_vals, v_path, Inches(0.5), v_top + Inches(0.35) + (i * (v_h + Inches(0.1))), width=Inches(9.0), height=v_h)
    add_image_bar(s_vals, "Valves Section View", Inches(0.5), v_top + Inches(5.0), Inches(9.0))

    # Safety / Gauge / Vacuum (3 Columns)
    s_misc = add_tank_slide(prs, cond_header, tank.tank_number, report_no, curr_page, True, False)
    curr_page += 1
    safe_v = insp_image_map.get("safetyvalve", [None])[0]
    gauge_v = insp_image_map.get("levelpressuregauge", [None])[0]
    vacuum_v = insp_image_map.get("vacuumreading", [None])[0]
    add_image_with_bars(s_misc, safe_v, Inches(0.3), Inches(1.15), Inches(3.0), Inches(4.0), "Safety Valve")
    add_image_with_bars(s_misc, gauge_v, Inches(3.5), Inches(1.15), Inches(3.0), Inches(4.0), "Level / Pressure Gauge")
    add_image_with_bars(s_misc, vacuum_v, Inches(6.7), Inches(1.15), Inches(3.0), Inches(4.0), "Vacuum Reading")

    # -------------------------------
    # END OF TANK CONDITION SECTION
    # -------------------------------
    insp_image_map = {}  # HARD RESET to ensure no bleed-over into subsequent slides
    print("DEBUG: insp_image_map reset. Keys:", list(insp_image_map.keys()))




    # --- OTHER SLIDES (Original functionality maintained) ---
    page_start = 10

    # Collect paths already used in Tank Condition slides (Slides 4-11) AND ALL NEW IMAGE SOURCES
    seen_paths = set()

    # 1. Add inspection images (tank condition slides)
    for paths in insp_image_map.values():
        for p in paths:
            if p:
                key = get_image_key(p, tank.tank_number, base_dir)
                seen_paths.add(key)

                resolved = resolve_path(p, tank.tank_number, base_dir)
                if resolved:
                    resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                    seen_paths.add(resolved_key)

    # 2. Add tank details images
    if d.tank_number_image_path:
        key = get_image_key(d.tank_number_image_path, tank.tank_number, base_dir)
        seen_paths.add(key)
        resolved = resolve_path(d.tank_number_image_path, tank.tank_number, base_dir)
        if resolved:
            resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
            seen_paths.add(resolved_key)

    # 3. Add certificate images
    if certs:
        cert = certs[0]
        if hasattr(cert, "periodic_inspection_pdf_path") and cert.periodic_inspection_pdf_path:
            key = get_image_key(cert.periodic_inspection_pdf_path, tank.tank_number, base_dir)
            seen_paths.add(key)
            resolved = resolve_path(cert.periodic_inspection_pdf_path, tank.tank_number, base_dir)
            if resolved:
                resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                seen_paths.add(resolved_key)

        if hasattr(cert, "next_insp_pdf_path") and cert.next_insp_pdf_path:
            key = get_image_key(cert.next_insp_pdf_path, tank.tank_number, base_dir)
            seen_paths.add(key)
            resolved = resolve_path(cert.next_insp_pdf_path, tank.tank_number, base_dir)
            if resolved:
                resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                seen_paths.add(resolved_key)

    # 4. Add valve label images
    # 4 & 5. Add valve label & frame images
    if tank_frame_outer_record:
        # GA Drawing
        if tank_frame_outer_record.ga_image_path:
            key = get_image_key(tank_frame_outer_record.ga_image_path, tank.tank_number, base_dir)
            seen_paths.add(key)
            resolved = resolve_path(tank_frame_outer_record.ga_image_path, tank.tank_number, base_dir)
            if resolved:
                resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                seen_paths.add(resolved_key)

        # Image 2
        if tank_frame_outer_record.image2_image_path:
            key = get_image_key(tank_frame_outer_record.image2_image_path, tank.tank_number, base_dir)
            seen_paths.add(key)
            resolved = resolve_path(tank_frame_outer_record.image2_image_path, tank.tank_number, base_dir)
            if resolved:
                resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                seen_paths.add(resolved_key)

        # Extra Valve & Shell images (3-6)
        for i in range(3, 7):
            fld = f"img{i}_path"
            val = getattr(tank_frame_outer_record, fld, None)
            if val:
                key = get_image_key(val, tank.tank_number, base_dir)
                seen_paths.add(key)
                resolved = resolve_path(val, tank.tank_number, base_dir)
                if resolved:
                    resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                    seen_paths.add(resolved_key)

    # 6. Add other images
    for ot in others_rows:
        if ot.image_path:
            key = get_image_key(ot.image_path, tank.tank_number, base_dir)
            seen_paths.add(key)
            resolved = resolve_path(ot.image_path, tank.tank_number, base_dir)
            if resolved:
                resolved_key = get_image_key(resolved, tank.tank_number, base_dir)
                seen_paths.add(resolved_key)
    
    # --- SLIDES 12-14 REMOVED AS PER REQUEST ---
    # Start fresh tracking for remaining dynamic slides
    slides_12_18_seen = set()



    # --- GROUP 1: Tank Drawings (P&ID, GA, and Images 3-6) ---
    if drawings:
        dwg = drawings[0]
        # 1. P&ID Drawing
        if hasattr(dwg, "pid_drawing") and dwg.pid_drawing:
             curr_page = add_full_page_image_slide(
                prs, tank.tank_number, report_no, curr_page,
                "P&ID Drawings",                       # Slide Title
                "P&ID Drawings",                       # Top Bar
                f"P&ID Reference: {dwg.pid_reference or '-'}",  # Bottom Bar
                dwg.pid_drawing, base_dir, exclude_paths=slides_12_18_seen
            )
        
        # 2. Image 2 Drawing
        if hasattr(dwg, "image2_drawing_file") and dwg.image2_drawing_file:
             curr_page = add_full_page_image_slide(
                prs, tank.tank_number, report_no, curr_page,
                "P&ID Drawings",                       # Slide Title
                "Image 2",                             # Top Bar
                "",                                    # Bottom Bar
                dwg.image2_drawing_file, base_dir, exclude_paths=slides_12_18_seen
            )

        # 3. Extra Drawings Images (3-6)
        for i in range(3, 7):
            fld = f"img{i}"
            img_path = getattr(dwg, fld, None)
            if img_path:
                curr_page = add_full_page_image_slide(
                    prs, tank.tank_number, report_no, curr_page,
                    "P&ID Drawings",                       # Slide Title
                    f"Drawing Image {i}",                  # Top Bar
                    "",                                    # Bottom Bar
                    img_path, base_dir, exclude_paths=slides_12_18_seen
                )

    # --- GROUP 2: Valve & Shell (Valve Label, Frame, and Images 3-6) ---
    if tank_frame_outer_record:
        # 1. GA Drawing
        if tank_frame_outer_record.ga_image_path:
            curr_page = add_full_page_image_slide(
                prs, tank.tank_number, report_no, curr_page,
                "Tank Frame & Outer Shell",                           # Slide Title
                "Tank Frame & Outer Shell",                           # Top Bar
                f"GA Reference: {tank_frame_outer_record.ga_reference or '-'}", # Bottom Bar
                tank_frame_outer_record.ga_image_path, base_dir, exclude_paths=slides_12_18_seen
            )

        # 2. Image 2
        if tank_frame_outer_record.image2_image_path:
            curr_page = add_full_page_image_slide(
                prs, tank.tank_number, report_no, curr_page,
                "Tank Frame & Outer Shell",                         # Slide Title
                "Image 2",                             # Top Bar
                "",                                    # Bottom Bar
                tank_frame_outer_record.image2_image_path, base_dir, exclude_paths=slides_12_18_seen
            )

        # 3. Extra Valve & Shell Images (3-6)
        for i in range(3, 7):
            fld = f"img{i}_path"
            img_path = getattr(tank_frame_outer_record, fld, None)
            if img_path:
                curr_page = add_full_page_image_slide(
                    prs, tank.tank_number, report_no, curr_page,
                    "Tank Frame & Outer Shell",                         # Slide Title
                    f"Image {i}",            # Top Bar
                    "",                                    # Bottom Bar
                    img_path, base_dir, exclude_paths=slides_12_18_seen
                )

    # --- LAST SLIDE (Contact Sheet) REMOVED ---



    # Return PPT as BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
